from django.shortcuts import render
from django.http import HttpResponse
from django.conf import settings
from django.http import JsonResponse 
import imgkit
import os
import PyPDF2
import spacy
import re
import numpy as np
import joblib
from django.http import HttpResponse
from tensorflow.keras.models import load_model
from tensorflow.keras.preprocessing.sequence import pad_sequences
from sklearn.metrics.pairwise import cosine_similarity
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.text_rank import TextRankSummarizer
import fitz
from PIL import Image
import io
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import linear_kernel
from collections import Counter
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
from django.template import Template, Context
import requests
import serpapi
from io import BytesIO
from transformers import BartTokenizer, BartForConditionalGeneration


API_URL = "https://api-inference.huggingface.co/models/nerijs/pixel-art-xl"
headers = {"Authorization": "Bearer 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"}

API_URL1 = "https://api-inference.huggingface.co/models/ml6team/keyphrase-extraction-kbir-inspec"
headers1 = {"Authorization": "Bearer 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"}

def query(payload):
    response = requests.post(API_URL, headers=headers, json=payload)
    return response.content

def query2(payload):
	response = requests.post(API_URL1, headers=headers1, json=payload)
	return response.json()

# You can access the image with PIL.Image for example
import io
from PIL import Image

# Load the English language model in spaCy
nlp = spacy.load("en_core_web_sm")

# Define a regular expression pattern to split text into sentences
sentence_pattern = r"(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?|\!)\s"

# With these lines using the constants from settings.py
tokenizer = joblib.load(settings.TOKENIZER_PATH)
label_encoder = joblib.load(settings.LABEL_ENCODER_PATH)
model = load_model(settings.MODEL_PATH)
results = []

def index(request):
    return HttpResponse("Hello")

# Define a view for rendering the upload form
def upload_form(request):
    return render(request, 'upload_form.html')

def upload_file(request):
    if request.method == 'POST' and request.FILES['file']:
        uploaded_file = request.FILES['file']

        if uploaded_file:
            # Specify the upload directory
            upload_dir = os.path.join(settings.MEDIA_ROOT, 'uploads')
            os.makedirs(upload_dir, exist_ok=True)

            # Save the uploaded file
            file_path = os.path.join(upload_dir, uploaded_file.name)
            with open(file_path, 'wb') as destination:
                for chunk in uploaded_file.chunks():
                    destination.write(chunk)

        results.clear() 
        section_texts = process_pdf(file_path)
        selected_model = request.POST.get('layout_choice', 'default_layout')
        selected_layout = request.POST.get('layout_choice2', 'default_layout')

        #response_data = {
        #   str(section): text for section, text in section_texts.items()
        #}

        #return JsonResponse(response_data)
        
        # Perform extractive summarization using TextRankx
        # Initialize an empty dictionary to store section headings and text
        section_texts = {}

        # Iterate through the sentences and their classifications
        for sentence, classification in results:
            if classification not in section_texts:
                section_texts[classification] = []
            section_texts[classification].append(sentence)

        # Perform extractive summarization using TextRank
        for section, text in section_texts.items():

            # Combine section text into a single string
            section_text = " ".join(text)

            # Set the desired number of sentences for each section
            if section in [0]:
                num_sentences = 2  # 2-3 sentences for conclusion
                text = extractive_summarization(selected_model,section_text, num_sentences=num_sentences)
            elif section in [4]:
                num_sentences = 4  # 2-3 sentences for objective
                text = extractive_summarization(selected_model,section_text, num_sentences=num_sentences)
            elif section in [1]:
                num_sentences = 2  # 2-3 sentences for objective
                text = extractive_summarization(selected_model,section_text, num_sentences=num_sentences)
            else:
                num_sentences = 0  # Default to 3 sentences
                text = extractive_summarization(selected_model,section_text, num_sentences=num_sentences)

            enhanced_summary = enhance_and_connect_summaries({section: [text]})

            # Replace placeholders in your template with the summarized text
            if section in [0]:
                abstract_text = enhanced_summary[section]
            elif section in [4]:
                results_text = enhanced_summary[section]
            elif section in [1]:
                conclusion_text = enhanced_summary[section]
        image_paths, ext_title=extract_figures_and_captions(file_path, abstract_text,results_text,conclusion_text)
        create_presentation_with_layout(selected_layout, abstract_text, results_text, conclusion_text, image_paths, ext_title)
        #generate_poster_image(selected_layout,abstract_text, results_text, conclusion_text, image_paths)
        return render(request, 'download_screen.html')

    return JsonResponse({'error': 'No file selected'})

from fuzzywuzzy import process
# Function to extract text from a PDF
def process_pdf(pdf_file_path):
    # Read the PDF and extract text
    # Example usage:
    text = ""
    stop_extraction = False  # Flag to stop extraction when "References" is found
    with open(pdf_file_path, "rb") as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        num_pages = len(pdf_reader.pages)

        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            page_text = page.extract_text()
            
            # Check if "References" is in the page text (case-insensitive)
            if re.search(r"references", page_text, re.IGNORECASE):
                stop_extraction = True  # Set the flag to stop extraction
            
            if not stop_extraction:
                text += page_text
    doc = nlp(text)
    # Split the text into sentences using the pattern
    sentences = re.split(sentence_pattern, text)

    # Initialize an empty list to store classification results

    for sentence in sentences:
        # Preprocess the sentence (e.g., lowercase, remove punctuation)
        #sentence = sentence.lower()
        #sentence = sentence.translate(str.maketrans('', '', string.punctuation))
        
        # Tokenize and pad the sentence
        max_sequence_length = 100
        sentence_sequence = tokenizer.texts_to_sequences([sentence])
        sentence_sequence = pad_sequences(sentence_sequence, maxlen=max_sequence_length)
        
        # Classify the sentence
        prediction = model.predict(sentence_sequence)
        predicted_label = label_encoder.inverse_transform([np.argmax(prediction)])[0]
        
        results.append((sentence, predicted_label))

    # Initialize an empty dictionary to store text by section
    section_texts = {}
    # Iterate through the sentences and their classifications
    for sentence, classification in results:
        if classification not in section_texts:
            section_texts[classification] = []
        section_texts[classification].append(sentence)
    return section_texts

# Define a function for extractive summarization using TextRank
from transformers import pipeline

pipe = pipeline("summarization", model="Falconsai/medical_summarization")
# Define a function to enhance and connect TextRank-generated summaries
def enhance_and_connect_summaries(section_texts):
    enhanced_summaries = {}  # Dictionary to store enhanced summaries for each section

    # Define a list of transition phrases
    transition_phrases = ["Additionally,", "Furthermore,", "In contrast,", "Moreover,", "On the other hand,", "However,", "Therefore,", "Likewise,", "Hence,"]

    for section, text in section_texts.items():
        # Reorder and group sentences logically if needed
        # Add transition phrases between sentences
        enhanced_summary = []
        prev_sentence = None
        for sentence in text:
            if prev_sentence:
                # Add a transition phrase before the sentence
                transition_phrase = transition_phrases.pop(0)
                enhanced_summary.append(transition_phrase)

            enhanced_summary.append(sentence)
            prev_sentence = sentence

        # Join the sentences into a coherent summary
        summary = " ".join(enhanced_summary)

        enhanced_summaries[section] = summary

    return enhanced_summaries

def extractive_summarization(model,text, num_sentences):
    # Initialize TextRank summarizer
    table_pattern = re.compile(r'\btable\b', flags=re.IGNORECASE)
    if(model=="layout1"):
        summarizer = TextRankSummarizer()

        # Tokenize the text
        parser = PlaintextParser.from_string(text, Tokenizer("english"))

        # Get the summary
        summary = summarizer(parser.document, num_sentences)

        # Reverse the order of sentences
        reversed_summary = list(summary)

        # Remove the last sentence
        reversed_summary = reversed_summary[:-1]

        # Filter sentences containing "table" variations using regex
        filtered_summary = [str(sentence) for sentence in reversed_summary if not table_pattern.search(str(sentence))]

        return " ".join(filtered_summary)

    elif(model=="layout3"):
        max_length = num_sentences * 150
        min_length = max_length - 50 
        summary = pipe(text, max_length, min_length, do_sample=False)
        return summary[0]['summary_text']
    elif model == "layout4":
        if num_sentences == 0:
            return text
        model_name = "facebook/bart-large-cnn"
        tokenizer = BartTokenizer.from_pretrained(model_name)
        model = BartForConditionalGeneration.from_pretrained(model_name)

        def generate_summary(model, tokenizer, text, max_length):
            inputs = tokenizer(text, return_tensors="pt", truncation=True, max_length=max_length)
            output = model.generate(**inputs)
            summary = tokenizer.decode(output[0], skip_special_tokens=True)
            return summary

        summary = generate_summary(model, tokenizer, text, max_length=1024)  # Adjust max_length as needed

        # Filter sentences containing "table" variations using regex
        filtered_summary = [sentence for sentence in summary.split(". ") if not table_pattern.search(sentence)]

        # Join the filtered sentences up to the specified number of sentences
        final_summary = ". ".join(filtered_summary[:num_sentences])

        return final_summary
    return ''


import fitz  # PyMuPDF
import re

def extract_title_heuristic(pdf_path, keyphrases):
    doc = fitz.open(pdf_path)

    max_similarity = 0
    best_title = None

    # Create TF-IDF vectors for the keyphrases
    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform(keyphrases)

    for page_num in range(len(doc)):  # Iterate over all pages
        page = doc[page_num]
        for block in page.get_text("blocks"):
            # Calculate TF-IDF vector for the current block
            block_tfidf = vectorizer.transform([block[4]])

            # Calculate cosine similarity
            similarity_scores = cosine_similarity(block_tfidf, tfidf_matrix)

            # Find the maximum similarity score
            max_block_similarity = similarity_scores.max()

            if max_block_similarity > 0.7:
                title_text = block[4].replace('\n', '')
                best_title = title_text
                break  # Stop iterating over blocks if a suitable title is found

        if best_title:  # Exit loop if a title is found on any page
            break

    return best_title

def extract_figures_and_captions(pdf_path, abstract, results, conclusion):
    figures_dict = {}
    current_figure_number = None
    current_caption = ''

    doc = fitz.open(pdf_path)

    for page_num in range(doc.page_count):
        page = doc[page_num]
        text = page.get_text()

        for line in text.split('\n'):
            match = re.search(r'\b(?:figure|fig)\.?\s*(\d+)\.?\b', line, re.IGNORECASE)
            
            if match:
                if current_figure_number:
                    figures_dict[current_figure_number] = figures_dict.get(current_figure_number, '') + current_caption

                current_figure_number = f"Figure {match.group(1)}"
                current_caption = line[match.end():].strip()
            elif current_figure_number:
                current_caption += ' ' + line.strip()

    if current_figure_number:
        figures_dict[current_figure_number] = figures_dict.get(current_figure_number, '') + current_caption

    figures_and_captions= figures_dict
    images_with_captions = []

    # Iterate over PDF pages
    for page_index in range(len(doc)):
        # Get the page itself
        page = doc[page_index]

        # Get the images on the page
        images = page.get_images(full=True)

        # Printing the number of images found on this page
        if images:
            print(f"[+] Found a total of {len(images)} images in page {page_index}")
        else:
            print("[!] No images found on page", page_index)

        # Iterate over images in reverse order if there are more than two images
        for image_index, img in enumerate(reversed(images), start=1):
            # Get the XREF of the image
            xref = img[0]

            # Extract the image bytes
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]

            # Get the image extension
            image_ext = base_image["ext"]

            # Create a PIL Image from the bytes
            pil_image = Image.open(io.BytesIO(image_bytes))

            # Get the caption corresponding to the image
            # This assumes you have a dictionary of captions named 'figures_and_captions'
            figure_key = f"Figure {image_index}"
            caption = figures_and_captions.get(figure_key, f"Image {image_index}")

            # Append the image along with its caption to the list
            images_with_captions.append((pil_image, caption))

    doc.close()

    #image_bytes = query({"inputs": abstract})
    #image = Image.open(io.BytesIO(image_bytes))
    #image.save("generated_pixel_art.png")
    
    summary = abstract +results+conclusion

    output1 = query2({"inputs": summary,})

    # Sorting the list of dictionaries based on the 'score' key in descending order
    sorted_output = sorted(output1, key=lambda x: int(x['score']), reverse=True)

    sorted_output = sorted_output[:2]
    
    keyphrases=[]
    
    for i in sorted_output:
        word = i['word']
        keyphrases.append(word)
        my_query = word+"vector graphics png"
        params = {
            "engine": "google_images",
            "q": my_query,
            "api_key": "10a5386439f58339f52184c15f77fe0f4437bda355931308153cfbcb4ba44303"
        }

        '''
        search = serpapi.search(params)
        results = search.as_dict()

        # Check if the 'images_results' key exists in the results
        if 'images_results' in results:
            # Get the 'images_results' list
            images_results = results['images_results']

            # Initialize a list to store original image URLs
            original_urls = []

            # Iterate over each image result
            for image_result in images_results:
                # Check if we already collected 2 original image URLs
                if len(original_urls) == 1:
                    break  # Stop the iteration if we have collected 2 URLs

                # Check if the 'original' key is present in the image result
                if 'original' in image_result:
                    # If 'original' key is present, add its value to the list
                    original_urls.append(image_result['original'])

            # Print the list of original image URLs
            print(original_urls)
        else:
            print("No image results found.")

        response = requests.get(original_urls[0])
        
        # Check if the request was successful (status code 200)
        if response.status_code == 200:
            # Open the image using PIL's Image module
            img = Image.open(BytesIO(response.content))
            
            # Display the image
            img.show()
        else:
            print(f"Failed to retrieve image from URL: {original_urls[0]}")
    '''
    extracted_title = extract_title_heuristic(pdf_path, keyphrases)
    if extracted_title:
        print(f"Extracted title: {extracted_title}")
    else:
        print("Title not found using heuristics.")

    # Create a list to store images along with their captions and similarities
    images_with_captions_and_similarities = []
    top_image_paths=[]

    # Create a TF-IDF vectorizer
    vectorizer = TfidfVectorizer()

    # Fit and transform the summary
    summary_vector = vectorizer.fit_transform([summary])

    # Iterate over images
    for image_index, (pil_image, caption) in enumerate(images_with_captions, start=1):
        # Fit and transform the caption
        caption_vector = vectorizer.transform([caption])

        # Calculate cosine similarity
        similarity = cosine_similarity(summary_vector, caption_vector)[0][0]

        # Append the image along with its caption and similarity to the list
        images_with_captions_and_similarities.append((pil_image, caption, similarity,image_index))
    
    images_with_captions_and_similarities = sorted(images_with_captions_and_similarities, key=lambda x: x[2], reverse=True)

# Now sorted_images contains the tuples sorted by similarity in descending order

    # Check if the list is not empty
    if images_with_captions_and_similarities:
        # Display the top 2 most similar images
        for i in range(min(2, len(images_with_captions_and_similarities))):
            similar_image, similar_caption, _ , img_index = images_with_captions_and_similarities[i]
            # Save the image to a file
            image_path = f"image_{i + 1}.jpeg"
            similar_image.save(image_path)

            # Add the image path and caption to the list
            top_image_paths.append((image_path, similar_caption, img_index))

        return top_image_paths, extracted_title
    else:
        print("No images with captions found.")
        return [],extracted_title
from pptx.util import Pt
from pptx.dml.color import RGBColor  # Corrected import statement
from pptx.enum.text import PP_ALIGN

def create_presentation_with_layout(layout, abstract_text, results_text, conclusion_text, top_images, title):
    prs = Presentation('poster_2.pptx')

    # Update text elements in the slides
    updates = {
        "abstract": abstract_text,
        "results": results_text,
        "conclusion": conclusion_text,
        "title": str(title),
        # Add more names and corresponding text updates as needed
    }

    for element_id, new_text in updates.items():
        # Find the shape or placeholder with the specified name
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name == element_id:
                    paragraph = shape.text_frame.paragraphs[0]
                    run = paragraph.runs[0]
                    run.text = new_text.strip()  # Update with new text

                    if(element_id!='title'):
                        # Set text size and font style
                        font = run.font
                        font.size = Pt(40)  # Set text size
                        font.name = 'Open Sans'  # Set font style

                        # Set text color for "conclusion"
                        if element_id == "conclusion":
                            font.color.rgb = RGBColor(255, 255, 255)  # Set color to white
                        # Remove extra line at the start of each paragraph
                        paragraph.space_before = Pt(0)
                    else:
                        # Set font properties
                        font = run.font
                        font.bold = True
                        font.size = Pt(70)  # Set font size
                        font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white
                        
                        # Set font family
                        font.name = 'Open Sans Bold'
                        
                        # Set alignment to center
                        paragraph.alignment = PP_ALIGN.CENTER
                        
                    # Set line spacing to default (1)
                    paragraph.line_spacing = 1.0
                

     # Replace placeholders with images
    for i, (image_path, _ , img_index) in enumerate(top_images):
        print(img_index)
        if i < 2:  # Assuming there are only two placeholders for images
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.name == f"image{i+1}":
                        shape.text = ''  # Clear the shape's text
                        if(i==0):
                            slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, 15700000)
                        else:
                            slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, 7800000)
    # Save the presentation to a file
    prs.save("my_presentation.pptx")